"""Extrae contenido de presentaciones PowerPoint (.pptx)."""

from __future__ import annotations


def extract_pptx(file_path: str) -> str:
    """Extrae texto, notas, tablas e imágenes de un archivo PowerPoint.

    Args:
        file_path: Ruta al archivo .pptx

    Returns:
        Markdown estructurado con contenido de cada slide
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    sections: list[str] = [
        f"**Presentación:** {total_slides} diapositivas\n",
    ]

    for idx, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"#### Diapositiva {idx} de {total_slides}"]

        # Layout name
        if slide.slide_layout and slide.slide_layout.name:
            slide_parts.append(f"*Layout: {slide.slide_layout.name}*\n")

        # Extract text from all shapes
        texts: list[str] = []
        tables_md: list[str] = []
        image_count = 0

        for shape in slide.shapes:
            # Text frames (titles, bullets, text boxes)
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # Detect heading level by font size or placeholder type
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                            ph_type = shape.placeholder_format.type
                            # Title placeholders
                            if ph_type in (1, 15, 16):  # TITLE, SUBTITLE variants
                                texts.append(f"**{text}**")
                                continue
                        # Bullet detection
                        if para.level > 0:
                            indent = "  " * para.level
                            texts.append(f"{indent}- {text}")
                        else:
                            texts.append(text)

            # Tables
            if shape.has_table:
                table = shape.table
                rows_data: list[list[str]] = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    rows_data.append(row_data)

                if rows_data:
                    # First row as header
                    headers = rows_data[0]
                    md_table = "| " + " | ".join(headers) + " |\n"
                    md_table += "| " + " | ".join(["---"] * len(headers)) + " |\n"
                    for row in rows_data[1:]:
                        # Pad row to match headers
                        while len(row) < len(headers):
                            row.append("")
                        md_table += "| " + " | ".join(row[:len(headers)]) + " |\n"
                    tables_md.append(md_table)

            # Images (just count — can't easily extract to text)
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image_count += 1

        if texts:
            slide_parts.append("\n".join(texts))

        if tables_md:
            for tbl in tables_md:
                slide_parts.append(f"\n{tbl}")

        if image_count:
            slide_parts.append(f"\n*[{image_count} imagen(es) en esta diapositiva]*")

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_parts.append(f"\n> **Notas del presentador:** {notes_text}")

        sections.append("\n".join(slide_parts))

    return "\n\n---\n\n".join(sections)
